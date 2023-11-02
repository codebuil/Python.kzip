import os
from pptx import Presentation
from PIL import Image
#pip install python-pptx
#pip install Pillow
def criar_apresentacao_ppt(nome_saida):
    ii = 0
  
    

    while os.path.exists(f"{ii}.bmp"):
        ii += 1
    presentation = Presentation()
    slides = [f"{i}.bmp" for i in range(0, ii)]  # Substitua o intervalo pelo número de slides desejado
    print(ii)
    ii=0
    for slide in slides:
        slide_layout = presentation.slide_layouts[5]  # Escolha o layout adequado para os slides

        sslide = presentation.slides.add_slide(slide_layout)
        img = Image.open(slide)
        left = top = 0

        # Redimensiona a imagem para se ajustar ao slide (altere conforme necessário)
        slide_width = presentation.slide_width
        slide_height = presentation.slide_height
        img.thumbnail((slide_width, slide_height))

        sslide.shapes.add_picture(slide, left, top, width=slide_width, height=slide_height)

    presentation.save(nome_saida)

if __name__ == "__main__":
    print("\x1bc\x1b[43;30m")

    nome_arquivo_saida = "out.pptx"
    criar_apresentacao_ppt(nome_arquivo_saida)
